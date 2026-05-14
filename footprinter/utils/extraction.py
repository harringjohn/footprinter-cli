"""
Text extraction from document formats for MCP read tool.

Extracts plaintext from binary document data (.docx, .pdf, .xlsx, .pptx, .csv, .tsv).
Sanitizes output to only include visible text (no comments, tracked changes, formulas).
"""

import csv
import io
import logging
from typing import Optional, Tuple

logger = logging.getLogger(__name__)

# Map file extensions to extractor types
EXTENSION_MAP = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".csv": "csv",
    ".tsv": "tsv",
}

# Map MIME types to extractor types
MIME_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/csv": "csv",
    "text/tab-separated-values": "tsv",
    # Google Workspace types (handled as text after Drive export)
    "application/vnd.google-apps.document": "text",
    "application/vnd.google-apps.spreadsheet": "text",
    "application/vnd.google-apps.presentation": "text",
}


def get_extractor_for_file(name: str, mime_type: str = "") -> Optional[str]:
    """
    Determine the extractor type for a file.

    Args:
        name: Filename with extension
        mime_type: MIME type hint (optional)

    Returns:
        Extractor type string ('pdf', 'docx', etc.) or None if no extraction needed
    """
    # Check MIME type first (more reliable)
    if mime_type and mime_type in MIME_MAP:
        return MIME_MAP[mime_type]

    # Fall back to extension
    ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""
    return EXTENSION_MAP.get(ext)


def extract_text(data: bytes, extractor_type: str) -> Tuple[Optional[str], Optional[str]]:
    """
    Extract text from binary document data.

    Args:
        data: Raw file bytes
        extractor_type: Type of extractor to use ('pdf', 'docx', etc.)

    Returns:
        Tuple of (extracted_text, error_message)
        - On success: (text, None)
        - On failure: (None, error_description)
    """
    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
        "csv": _extract_csv,
        "tsv": _extract_tsv,
        "text": _extract_text,
    }

    extractor = extractors.get(extractor_type)
    if not extractor:
        return None, f"Unknown extractor type: {extractor_type}"

    try:
        text = extractor(data)
        return text, None
    except ImportError as e:
        return None, f"Missing dependency: {e}"
    except Exception as e:
        logger.error(f"Extraction error ({extractor_type}): {e}")
        return None, str(e)


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF bytes."""
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(data))
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)

    return "\n\n".join(text_parts)


def _extract_docx(data: bytes) -> str:
    """
    Extract text from DOCX bytes.

    Only extracts visible paragraph text - no comments, tracked changes, or headers/footers.
    """
    import docx

    doc = docx.Document(io.BytesIO(data))
    text_parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_parts.append(" | ".join(row_text))

    return "\n\n".join(text_parts)


def _extract_xlsx(data: bytes) -> str:
    """
    Extract text from XLSX bytes.

    Extracts cell values only - no formulas, comments, or hidden data.
    """
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"=== {sheet_name} ===")

        for row in sheet.iter_rows(values_only=True):
            # Filter out None values and convert to strings
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                text_parts.append(" | ".join(row_values))

    wb.close()
    return "\n".join(text_parts)


def _extract_pptx(data: bytes) -> str:
    """
    Extract text from PPTX bytes.

    Extracts text from shapes and text frames - no speaker notes or comments.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        text_parts.append(f"--- Slide {slide_num} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        if slide_text:
            text_parts.append("\n".join(slide_text))

    return "\n\n".join(text_parts)


def _extract_csv(data: bytes) -> str:
    """Extract text from CSV bytes."""
    # Try to decode as UTF-8 first, fall back to latin-1
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1")

    # Parse and format as readable text
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(row))

    return "\n".join(rows)


def _extract_tsv(data: bytes) -> str:
    """Extract text from TSV bytes."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1")

    reader = csv.reader(io.StringIO(text), delimiter="\t")
    rows = []
    for row in reader:
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(row))

    return "\n".join(rows)


def _extract_text(data: bytes) -> str:
    """Pass-through for already-text content (e.g., Google Workspace exports)."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")
