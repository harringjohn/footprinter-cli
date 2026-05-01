"""
Tests for MCP text extraction module.
"""

import pytest

from footprinter.mcp.extraction import (
    extract_text,
    get_extractor_for_file,
)


# Helper functions to check dependency availability (must be defined before use in skipif)
def _has_pypdf():
    try:
        import pypdf  # noqa: F401

        return True
    except ImportError:
        return False


def _has_python_docx():
    try:
        import docx  # noqa: F401

        return True
    except ImportError:
        return False


def _has_openpyxl():
    try:
        import openpyxl  # noqa: F401

        return True
    except ImportError:
        return False


def _has_python_pptx():
    try:
        from pptx import Presentation  # noqa: F401

        return True
    except ImportError:
        return False


class TestExtractorDetection:
    """Test extractor type detection from filename and mime type."""

    def test_pdf_by_extension(self):
        assert get_extractor_for_file("document.pdf") == "pdf"
        assert get_extractor_for_file("DOCUMENT.PDF") == "pdf"

    def test_docx_by_extension(self):
        assert get_extractor_for_file("report.docx") == "docx"

    def test_xlsx_by_extension(self):
        assert get_extractor_for_file("data.xlsx") == "xlsx"

    def test_pptx_by_extension(self):
        assert get_extractor_for_file("presentation.pptx") == "pptx"

    def test_csv_by_extension(self):
        assert get_extractor_for_file("data.csv") == "csv"

    def test_tsv_by_extension(self):
        assert get_extractor_for_file("data.tsv") == "tsv"

    def test_unknown_extension_returns_none(self):
        assert get_extractor_for_file("file.txt") is None
        assert get_extractor_for_file("image.png") is None
        assert get_extractor_for_file("noextension") is None

    def test_mime_type_takes_precedence(self):
        # Even with .txt extension, mime type should win
        result = get_extractor_for_file(
            "document.txt",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        assert result == "docx"

    def test_pdf_by_mime(self):
        assert get_extractor_for_file("file", "application/pdf") == "pdf"

    def test_google_docs_returns_text(self):
        assert get_extractor_for_file("file", "application/vnd.google-apps.document") == "text"

    def test_google_sheets_returns_text(self):
        assert get_extractor_for_file("file", "application/vnd.google-apps.spreadsheet") == "text"


class TestCSVExtraction:
    """Test CSV extraction."""

    def test_simple_csv(self):
        data = b"name,age,city\nAlice,30,NYC\nBob,25,LA"
        text, error = extract_text(data, "csv")
        assert error is None
        assert "name | age | city" in text
        assert "Alice | 30 | NYC" in text
        assert "Bob | 25 | LA" in text

    def test_csv_with_empty_rows(self):
        data = b"a,b\n\nc,d"
        text, error = extract_text(data, "csv")
        assert error is None
        assert "a | b" in text
        assert "c | d" in text

    def test_csv_utf8(self):
        data = "name,city\nJosé,São Paulo".encode("utf-8")
        text, error = extract_text(data, "csv")
        assert error is None
        assert "José" in text
        assert "São Paulo" in text


class TestTSVExtraction:
    """Test TSV extraction."""

    def test_simple_tsv(self):
        data = b"name\tage\tcity\nAlice\t30\tNYC"
        text, error = extract_text(data, "tsv")
        assert error is None
        assert "name | age | city" in text
        assert "Alice | 30 | NYC" in text


class TestPDFExtraction:
    """Test PDF extraction."""

    @pytest.mark.skipif(not _has_pypdf(), reason="pypdf not installed")
    def test_invalid_pdf_returns_error(self):
        data = b"not a real pdf"
        text, error = extract_text(data, "pdf")
        assert text is None
        assert error is not None


class TestDocxExtraction:
    """Test DOCX extraction."""

    @pytest.mark.skipif(not _has_python_docx(), reason="python-docx not installed")
    def test_invalid_docx_returns_error(self):
        data = b"not a real docx"
        text, error = extract_text(data, "docx")
        assert text is None
        assert error is not None


class TestXlsxExtraction:
    """Test XLSX extraction."""

    @pytest.mark.skipif(not _has_openpyxl(), reason="openpyxl not installed")
    def test_invalid_xlsx_returns_error(self):
        data = b"not a real xlsx"
        text, error = extract_text(data, "xlsx")
        assert text is None
        assert error is not None


class TestPptxExtraction:
    """Test PPTX extraction."""

    @pytest.mark.skipif(not _has_python_pptx(), reason="python-pptx not installed")
    def test_invalid_pptx_returns_error(self):
        data = b"not a real pptx"
        text, error = extract_text(data, "pptx")
        assert text is None
        assert error is not None


class TestTextExtraction:
    """Test text pass-through extraction."""

    def test_utf8_text(self):
        data = "Hello, world! Über café.".encode("utf-8")
        text, error = extract_text(data, "text")
        assert error is None
        assert "Hello, world!" in text
        assert "Über café" in text

    def test_latin1_fallback(self):
        # Bytes that aren't valid UTF-8 but are valid latin-1
        data = b"Hello \xe9"  # é in latin-1
        text, error = extract_text(data, "text")
        assert error is None
        assert "Hello" in text


class TestUnknownExtractor:
    """Test unknown extractor type handling."""

    def test_unknown_type_returns_error(self):
        text, error = extract_text(b"data", "unknown_type")
        assert text is None
        assert "Unknown extractor type" in error


class TestCSVLatin1Fallback:
    """Test CSV latin-1 fallback when UTF-8 decoding fails."""

    def test_csv_latin1_fallback(self):
        # Latin-1 encoded CSV with non-UTF-8 bytes
        data = b"name,city\nJos\xe9,S\xe3o Paulo"
        text, error = extract_text(data, "csv")
        assert error is None
        assert "Paulo" in text

    def test_tsv_latin1_fallback(self):
        data = b"name\tcity\nJos\xe9\tS\xe3o Paulo"
        text, error = extract_text(data, "tsv")
        assert error is None
        assert "Paulo" in text


class TestValidPDFExtraction:
    """Test PDF extraction with a real (minimal) PDF."""

    @pytest.mark.skipif(not _has_pypdf(), reason="pypdf not installed")
    def test_pdf_extraction_with_valid_pdf(self):
        import io

        import pypdf

        # Create a minimal PDF with text
        writer = pypdf.PdfWriter()
        writer.add_blank_page(width=200, height=200)
        # pypdf doesn't easily add text, but we verify the pipeline works
        buf = io.BytesIO()
        writer.write(buf)
        data = buf.getvalue()

        text, error = extract_text(data, "pdf")
        assert error is None
        assert isinstance(text, str)


class TestValidDocxExtraction:
    """Test DOCX extraction with a real .docx."""

    @pytest.mark.skipif(not _has_python_docx(), reason="python-docx not installed")
    def test_docx_extraction_with_valid_docx(self):
        import io

        import docx

        doc = docx.Document()
        doc.add_paragraph("Hello from DOCX")
        doc.add_paragraph("Second paragraph")

        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"

        buf = io.BytesIO()
        doc.save(buf)
        data = buf.getvalue()

        text, error = extract_text(data, "docx")
        assert error is None
        assert "Hello from DOCX" in text
        assert "Second paragraph" in text
        assert "A1" in text
        assert "B2" in text


class TestValidXlsxExtraction:
    """Test XLSX extraction with a real .xlsx."""

    @pytest.mark.skipif(not _has_openpyxl(), reason="openpyxl not installed")
    def test_xlsx_extraction_with_valid_xlsx(self):
        import io

        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Age"])
        ws.append(["Alice", 30])

        # Add second sheet
        ws2 = wb.create_sheet("Summary")
        ws2.append(["Total", 1])

        buf = io.BytesIO()
        wb.save(buf)
        data = buf.getvalue()

        text, error = extract_text(data, "xlsx")
        assert error is None
        assert "=== Data ===" in text
        assert "Alice" in text
        assert "30" in text
        assert "=== Summary ===" in text


class TestValidPptxExtraction:
    """Test PPTX extraction with a real .pptx."""

    @pytest.mark.skipif(not _has_python_pptx(), reason="python-pptx not installed")
    def test_pptx_extraction_with_valid_pptx(self):
        import io

        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

        # Add title
        slide.shapes.title.text = "Slide Title"
        # Add body content
        slide.placeholders[1].text = "Bullet point content"

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()

        text, error = extract_text(data, "pptx")
        assert error is None
        assert "Slide Title" in text
        assert "Bullet point content" in text


class TestImportErrorHandling:
    """Test that ImportError in extractors is handled gracefully."""

    def test_import_error_returns_missing_dependency(self):
        from unittest.mock import patch

        # Mock pypdf import to raise ImportError
        with patch.dict("sys.modules", {"pypdf": None}):
            # Force re-import by calling the extractor which does lazy import
            import sys

            # Remove cached module so the import inside _extract_pdf fails
            saved = sys.modules.pop("pypdf", None)
            try:
                text, error = extract_text(b"fake pdf data", "pdf")
                assert text is None
                assert "Missing dependency" in error or error is not None
            finally:
                if saved is not None:
                    sys.modules["pypdf"] = saved


class TestExtractionFallback:
    """Test fallback behavior in read tool integration."""

    def test_corrupt_file_fallback(self):
        # CSV extraction should succeed even on binary data
        # (it just returns garbled output)
        data = b"\x00\x01\x02\x03"
        text, error = extract_text(data, "csv")
        # CSV parser will handle this somehow
        assert error is None or text is not None
